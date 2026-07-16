"""
Générateur de rapport PowerPoint - Phase 1
Lit ProjectReportPivot.json et produit un rapport .pptx.
Aucun calcul ici : toutes les valeurs viennent déjà calculées du pivot (Pandas).

Installation : pip install python-pptx
Utilisation  : python pptx_generator.py
"""

import json
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from report_constants import (
    STATUS_LABEL, PRIORITY_LABEL, STATUS_HEX_COLOR,
    HEX_NAVY, HEX_TEAL, HEX_GREEN, HEX_CORAL, HEX_GRAY, HEX_LIGHT,
)


def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Palette cohérente avec le reste du projet (dashboard, diagrammes)
NAVY     = _hex_to_rgb(HEX_NAVY)
TEAL     = _hex_to_rgb(HEX_TEAL)
GREEN    = _hex_to_rgb(HEX_GREEN)
CORAL    = _hex_to_rgb(HEX_CORAL)
GRAY     = _hex_to_rgb(HEX_GRAY)
LIGHT_BG = _hex_to_rgb(HEX_LIGHT)

STATUS_COLOR = {k: _hex_to_rgb(v) for k, v in STATUS_HEX_COLOR.items()}


def charger_pivot(chemin: str = "ProjectReportPivot.json") -> dict:
    with open(chemin, "r", encoding="utf-8") as f:
        return json.load(f)


def ajouter_slide_titre(prs: Presentation, pivot: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout vierge
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = pivot["project"]["nom_projet"]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    genere_le = pivot["reporting_period"]["genere_le"][:10]
    run2.text = f"Rapport généré le {genere_le}  —  Projet {pivot['project']['cle_jira']}"
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)


def ajouter_slide_kpis(prs: Presentation, pivot: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Indicateurs clés"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NAVY

    kpis = pivot.get("kpis", [])
    cols = 3
    card_w = Inches(2.9)
    card_h = Inches(1.6)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    start_y = Inches(1.5)

    for i, kpi in enumerate(kpis):
        row, col = divmod(i, cols)
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        tf2 = card.text_frame
        tf2.margin_left = Inches(0.15)
        tf2.margin_top = Inches(0.12)
        tf2.word_wrap = True

        p_val = tf2.paragraphs[0]
        r_val = p_val.add_run()
        unite = kpi.get("unite") or ""
        r_val.text = f"{kpi['valeur']:g}{unite}"
        r_val.font.size = Pt(26)
        r_val.font.bold = True
        r_val.font.color.rgb = TEAL

        p_label = tf2.add_paragraph()
        r_label = p_label.add_run()
        r_label.text = kpi["label"]
        r_label.font.size = Pt(11)
        r_label.font.color.rgb = GRAY


def ajouter_slide_taches(prs: Presentation, pivot: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Détail des tâches"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NAVY

    tasks = pivot.get("tasks", [])
    n_rows = len(tasks) + 1
    n_cols = 4

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.4 * n_rows)
    )
    table = table_shape.table

    headers = ["Ticket", "Titre", "Statut", "Priorité"]
    widths = [Inches(1.0), Inches(4.2), Inches(1.8), Inches(1.8)]
    for i, w in enumerate(widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(12)

    for r, task in enumerate(tasks, start=1):
        statut = task["statut"]
        valeurs = [
            task["task_id"],
            task["titre"],
            STATUS_LABEL.get(statut, statut),
            PRIORITY_LABEL.get(task["priorite"], task["priorite"]),
        ]
        for c, val in enumerate(valeurs):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 else LIGHT_BG
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(11)
            if c == 2:
                para.runs[0].font.color.rgb = STATUS_COLOR.get(statut, GRAY)
                para.runs[0].font.bold = True


def generer_rapport(chemin_pivot: str = "ProjectReportPivot.json", sortie: str = "rapport_demo.pptx"):
    pivot = charger_pivot(chemin_pivot)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    ajouter_slide_titre(prs, pivot)
    ajouter_slide_kpis(prs, pivot)
    ajouter_slide_taches(prs, pivot)

    prs.save(sortie)
    print(f"Rapport PPTX généré : {sortie}")


if __name__ == "__main__":
    generer_rapport()
